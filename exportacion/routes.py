# este archivo genera los documentos descargables del proyecto en formato pdf
# para entregar al profesor y en formato powerpoint (.pptx) para las diapositivas.
# su propósito es permitir que los estudiantes descarguen una copia imprimible o
# editable de su trabajo escolar. lo hace leyendo los datos de la base de datos
# y convirtiendo los títulos, listas de tareas y trazos visuales en archivos.
# elegimos la librería reportlab para los documentos pdf porque permite armar
# informes en hojas de tamaño carta listos para imprimir, y elegimos la librería
# python-pptx porque convierte las diapositivas web en archivos editables de
# powerpoint compatibles con cualquier computadora o proyector de la feria.
# se diseñó así para poder entregar informes físicos o proyectar exposiciones.


# importamos herramientas para codificación, manejo de memoria, texto y formatos xml
import base64

import binascii
import io
import json
import re
from xml.sax.saxutils import escape

# importamos las funciones de flask para descargas de archivos y manejo de errores
from flask import Blueprint, send_file, abort
# importamos la protección de sesión del usuario
from flask_login import login_required, current_user

# importamos la base de datos
from extensiones import db
# importamos los modelos de proyectos y títulos de estado de tareas
from proyectos.models import Proyecto
from tareas.models import TITULOS_ESTADO

# importamos el generador de documentos pdf (reportlab)
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

# importamos el generador de presentaciones en formato powerpoint (python-pptx)
from pptx import Presentation as ArchivoPptx
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

# agrupamos las opciones de descarga e impresión del proyecto bajo la sección exportación
exportacion = Blueprint("exportacion", __name__)

# constante para adaptar las medidas de pantalla a las medidas oficiales de powerpoint
EMU_POR_PIXEL = 9525
ANCHO_LIENZO = 960
ALTO_LIENZO = 540


# función auxiliar que verifica si el proyecto existe y si el alumno conectado es integrante del grupo.
# sirve para asegurar que solo los dueños del trabajo escolar puedan descargar el informe en pdf o la presentación.
def obtener_proyecto_de_miembro(id_proyecto):
    proyecto = db.get_or_404(Proyecto, id_proyecto)
    if not proyecto.equipo.es_miembro(current_user):
        abort(403)
    return proyecto


# función auxiliar que limpia el nombre del proyecto para convertirlo en un nombre de archivo seguro de descargar.
# sirve para evitar que símbolos o caracteres extraños dañen el archivo descargado en la computadora del alumno.
# lo hace quitando acentos y espacios vacíos para generar nombres limpios como feria_tecnologica.pdf.
def nombre_de_archivo(texto, extension):
    limpio = re.sub(r"[^\w\s-]", "", texto).strip().replace(" ", "_") or "vura"
    return f"{limpio}.{extension}"


# función auxiliar que extrae los textos y notas pegajosas dibujadas en la pizarra y los convierte en viñetas para el informe pdf.
# sirve para incluir la lluvia de ideas y bocetos dentro del reporte escrito impreso que se entrega al profesor.
def parrafos_desde_pizarra(contenido_json, estilos):
    try:
        objetos = json.loads(contenido_json).get("objetos", {})
    except (ValueError, AttributeError):
        return []

    parrafos = []
    for objeto in objetos.values():
        texto = (objeto.get("text") or "").strip()
        if objeto.get("type") in ("textbox", "i-text", "text") and texto:
            parrafos.append(Paragraph("• " + escape(texto.replace("\n", " ")), estilos["BodyText"]))
    return parrafos


# función auxiliar que convierte el contenido de texto almacenado en párrafos formateados para el informe impreso.
# sirve para conservar las viñetas, negrillas y estilos de texto al momento de construir el documento en pdf.
def parrafos_desde_delta(delta_json, estilos):
    try:
        operaciones = json.loads(delta_json).get("ops", [])
    except (ValueError, AttributeError):
        return []

    parrafos = []
    linea = ""
    for operacion in operaciones:
        texto = operacion.get("insert")
        if not isinstance(texto, str):
            continue
        atributos = operacion.get("attributes") or {}
        partes = texto.split("\n")
        for posicion, parte in enumerate(partes):
            if parte:
                fragmento = escape(parte)
                if atributos.get("bold"):
                    fragmento = f"<b>{fragmento}</b>"
                if atributos.get("italic"):
                    fragmento = f"<i>{fragmento}</i>"
                linea += fragmento
            if posicion < len(partes) - 1:
                estilo = estilos["Heading2"] if atributos.get("header") else estilos["BodyText"]
                parrafos.append(Paragraph(linea or " ", estilo))
                linea = ""
    if linea:
        parrafos.append(Paragraph(linea, estilos["BodyText"]))
    return parrafos


# función que genera y entrega un documento PDF completo con el resumen del proyecto, tareas asignadas y notas de la pizarra.
# sirve para que los estudiantes impriman el informe oficial de su trabajo escolar y se lo entreguen al docente evaluador.
# lo hace organizando el documento en hojas de tamaño carta con títulos, tablas de tareas y listas explicativas.
# se construyó de esta manera para cumplir con los requisitos de entrega escrita impresos de la feria tecnológica.

@exportacion.route("/proyectos/<int:id_proyecto>/exportar/pdf")
@login_required
def exportar_pdf(id_proyecto):
    proyecto = obtener_proyecto_de_miembro(id_proyecto)
    estilos = getSampleStyleSheet()

    # encabezado del documento pdf con nombre del proyecto y equipo
    contenido = [
        Paragraph(escape(proyecto.nombre), estilos["Title"]),
        Paragraph(f"Equipo: {escape(proyecto.equipo.nombre)}", estilos["Normal"]),
    ]
    if proyecto.fecha_entrega:
        contenido.append(Paragraph(
            f"Fecha de entrega: {proyecto.fecha_entrega.strftime('%d/%m/%Y')}", estilos["Normal"]))
    if proyecto.descripcion:
        contenido += [Spacer(1, 0.4 * cm), Paragraph(escape(proyecto.descripcion), estilos["BodyText"])]

    # sección de tareas asignadas en el informe pdf
    contenido += [Spacer(1, 0.6 * cm), Paragraph("Tareas", estilos["Heading1"])]
    if proyecto.tareas:
        for tarea in proyecto.tareas:
            detalles = [TITULOS_ESTADO[tarea.estado]]
            if tarea.asignado:
                detalles.append(tarea.asignado.nombre)
            if tarea.fecha_limite:
                detalles.append("entrega " + tarea.fecha_limite.strftime("%d/%m/%Y"))
            contenido.append(Paragraph(
                f"• <b>{escape(tarea.titulo)}</b> — {escape(', '.join(detalles))}", estilos["BodyText"]))
    else:
        contenido.append(Paragraph("Sin tareas registradas.", estilos["BodyText"]))

    # sección de notas y pizarra en el informe pdf
    contenido += [Spacer(1, 0.6 * cm), Paragraph("Pizarra", estilos["Heading1"])]
    contenido_nota = proyecto.nota.contenido_json if proyecto.nota else None
    parrafos_notas = (parrafos_desde_pizarra(contenido_nota, estilos)
                      or parrafos_desde_delta(contenido_nota, estilos))
    contenido += parrafos_notas or [Paragraph("La pizarra está vacía.", estilos["BodyText"])]

    archivo = io.BytesIO()
    # construye y entrega el archivo pdf descargable
    SimpleDocTemplate(archivo, pagesize=letter, title=proyecto.nombre).build(contenido)
    archivo.seek(0)

    return send_file(archivo, mimetype="application/pdf", as_attachment=True,
                     download_name=nombre_de_archivo(proyecto.nombre, "pdf"))


# función auxiliar que traduce códigos de color (como #FF0000) al formato numérico necesario para powerpoint.
# sirve para mantener exactamente los mismos colores elegidos por el alumno al exportar sus diapositivas.
def color_desde_hex(texto):
    if isinstance(texto, str) and re.fullmatch(r"#[0-9a-fA-F]{6}", texto):
        return RGBColor.from_string(texto[1:])
    return RGBColor(0, 0, 0)


ALINEACIONES = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


# función auxiliar que convierte una foto o imagen pegada en la pantalla web en un archivo apto para incrustar en powerpoint.
# sirve para preservar las imágenes de la exposición sin romper la presentación al descargar.
def imagen_desde_data_url(src):
    if not isinstance(src, str) or not src.startswith("data:image"):
        return None
    try:
        _, datos = src.split(",", 1)
        return io.BytesIO(base64.b64decode(datos))
    except (ValueError, binascii.Error):
        return None


# función auxiliar que toma cada elemento (casillas de texto, fotos, rectángulos, círculos) y lo dibuja dentro de la diapositiva de powerpoint.
# sirve para recrear fielmente en el archivo descargable de powerpoint todo el trabajo de diseño visual hecho en vura.
def agregar_objeto_fabric(diapositiva_pptx, objeto):
    tipo = objeto.get("type", "")
    izquierda = Emu(int(objeto.get("left", 0) * EMU_POR_PIXEL))
    arriba = Emu(int(objeto.get("top", 0) * EMU_POR_PIXEL))
    escala_x = objeto.get("scaleX", 1)
    escala_y = objeto.get("scaleY", 1)

    # si el elemento es una casilla de texto
    if tipo in ("i-text", "textbox", "text"):
        ancho = Emu(int(max(objeto.get("width", 200), 10) * escala_x * EMU_POR_PIXEL))
        alto = Emu(int(max(objeto.get("height", 50), 10) * escala_y * EMU_POR_PIXEL))
        caja = diapositiva_pptx.shapes.add_textbox(izquierda, arriba, ancho, alto)
        marco = caja.text_frame
        marco.word_wrap = True
        lineas = (objeto.get("text", "") or "").split("\n")
        for indice, linea in enumerate(lineas):
            parrafo = marco.paragraphs[0] if indice == 0 else marco.add_paragraph()
            parrafo.text = linea
            parrafo.alignment = ALINEACIONES.get(objeto.get("textAlign"), PP_ALIGN.LEFT)
            fuente = parrafo.font
            fuente.size = Pt(objeto.get("fontSize", 32) * escala_y * 0.75)
            fuente.bold = objeto.get("fontWeight") == "bold"
            fuente.italic = objeto.get("fontStyle") == "italic"
            fuente.color.rgb = color_desde_hex(objeto.get("fill"))
            nombre_fuente = objeto.get("fontFamily")
            if isinstance(nombre_fuente, str):
                fuente.name = nombre_fuente
        return

    # si el elemento es una imagen
    if tipo == "image":
        datos_imagen = imagen_desde_data_url(objeto.get("src"))
        if datos_imagen is None:
            return
        ancho = Emu(int(objeto.get("width", 100) * escala_x * EMU_POR_PIXEL))
        alto = Emu(int(objeto.get("height", 100) * escala_y * EMU_POR_PIXEL))
        try:
            diapositiva_pptx.shapes.add_picture(datos_imagen, izquierda, arriba, ancho, alto)
        except Exception:
            pass
        return

    # si el elemento es una línea o flecha
    if tipo == "line":
        ancho = int(objeto.get("width", 100) * escala_x * EMU_POR_PIXEL)
        alto = int(objeto.get("height", 0) * escala_y * EMU_POR_PIXEL)
        conector = diapositiva_pptx.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, izquierda, arriba,
            Emu(int(objeto.get("left", 0) * EMU_POR_PIXEL) + ancho),
            Emu(int(objeto.get("top", 0) * EMU_POR_PIXEL) + alto),
        )
        conector.line.color.rgb = color_desde_hex(objeto.get("stroke"))
        conector.line.width = Pt(max(objeto.get("strokeWidth", 3), 1))
        return

    # si el elemento es un grupo de figuras unidas
    if tipo == "group":
        for parte in objeto.get("objects", []):
            hijo = dict(parte)
            hijo["left"] = objeto.get("left", 0) + parte.get("left", 0) * escala_x
            hijo["top"] = objeto.get("top", 0) + parte.get("top", 0) * escala_y
            hijo["scaleX"] = parte.get("scaleX", 1) * escala_x
            hijo["scaleY"] = parte.get("scaleY", 1) * escala_y
            agregar_objeto_fabric(diapositiva_pptx, hijo)
        return

    # si el elemento es un rectángulo, círculo o triángulo
    if tipo == "rect":
        forma = MSO_SHAPE.RECTANGLE
        ancho = Emu(int(objeto.get("width", 100) * escala_x * EMU_POR_PIXEL))
        alto = Emu(int(objeto.get("height", 100) * escala_y * EMU_POR_PIXEL))
    elif tipo == "circle":
        forma = MSO_SHAPE.OVAL
        diametro = objeto.get("radius", 50) * 2
        ancho = Emu(int(diametro * escala_x * EMU_POR_PIXEL))
        alto = Emu(int(diametro * escala_y * EMU_POR_PIXEL))
    elif tipo == "triangle":
        forma = MSO_SHAPE.ISOSCELES_TRIANGLE
        ancho = Emu(int(objeto.get("width", 100) * escala_x * EMU_POR_PIXEL))
        alto = Emu(int(objeto.get("height", 100) * escala_y * EMU_POR_PIXEL))
    else:
        return

    figura = diapositiva_pptx.shapes.add_shape(forma, izquierda, arriba, ancho, alto)
    figura.fill.solid()
    figura.fill.fore_color.rgb = color_desde_hex(objeto.get("fill"))
    figura.line.fill.background()


# función que convierte la presentación de diapositivas del proyecto en un archivo real de PowerPoint (.pptx) descargable.
# sirve para que los estudiantes proyecten sus láminas durante la exposición o las abran en PowerPoint para editarlas más si lo desean.
# lo hace creando las láminas, agregando los fondos, figuras y textos, y enviando el archivo listo a la computadora del alumno.
# se construyó de este modo para garantizar la compatibilidad con proyectores y computadoras de la feria tecnológica.
@exportacion.route("/proyectos/<int:id_proyecto>/exportar/pptx")
@login_required
def exportar_pptx(id_proyecto):
    proyecto = obtener_proyecto_de_miembro(id_proyecto)
    if proyecto.presentacion is None:
        abort(404)

    archivo_pptx = ArchivoPptx()
    archivo_pptx.slide_width = Emu(ANCHO_LIENZO * EMU_POR_PIXEL)
    archivo_pptx.slide_height = Emu(ALTO_LIENZO * EMU_POR_PIXEL)
    diseno_vacio = archivo_pptx.slide_layouts[6]

    # convierte cada diapositiva creada en vura en una diapositiva real de powerpoint
    for diapositiva in proyecto.presentacion.diapositivas:
        diapositiva_pptx = archivo_pptx.slides.add_slide(diseno_vacio)
        if not diapositiva.contenido_json:
            continue
        try:
            datos = json.loads(diapositiva.contenido_json)
        except ValueError:
            continue
        fondo = datos.get("background")
        color_fondo = fondo if isinstance(fondo, str) else (fondo or {}).get("color")
        if color_fondo and color_fondo != "#ffffff":
            relleno = diapositiva_pptx.background.fill
            relleno.solid()
            relleno.fore_color.rgb = color_desde_hex(color_fondo)
        for objeto in datos.get("objects", []):
            agregar_objeto_fabric(diapositiva_pptx, objeto)

    archivo = io.BytesIO()
    archivo_pptx.save(archivo)
    archivo.seek(0)

    # entrega el archivo de powerpoint listo para descargar en la computadora del alumno
    return send_file(
        archivo,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=nombre_de_archivo(proyecto.nombre, "pptx"),
    )


